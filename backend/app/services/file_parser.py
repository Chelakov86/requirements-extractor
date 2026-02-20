import io
import os
from dataclasses import dataclass

import fitz  # PyMuPDF
from docx import Document as DocxDocument
from fastapi import HTTPException, UploadFile
from openpyxl import load_workbook
from pptx import Presentation

from app.config import settings


@dataclass
class ParsedDocument:
    filename: str
    file_type: str  # pdf | docx | xlsx | pptx | txt | md | text
    raw_text: str
    char_count: int


class FileParser:
    SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".txt", ".md", ".xlsx", ".pptx"}
    SUPPORTED_MIME_TYPES = {
        "application/pdf": ".pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
        "text/plain": ".txt",
        "text/markdown": ".md",
    }

    def validate_file(self, filename: str, content_type: str, size_bytes: int) -> None:
        ext = os.path.splitext(filename)[1].lower()

        if ext not in self.SUPPORTED_EXTENSIONS:
            raise HTTPException(
                status_code=400,
                detail={
                    "error": "INVALID_FILE_TYPE",
                    "message": f"File '{filename}' has unsupported extension '{ext}'. "
                    f"Supported: {', '.join(sorted(self.SUPPORTED_EXTENSIONS))}",
                },
            )

        expected_ext = self.SUPPORTED_MIME_TYPES.get(content_type)
        if expected_ext is None or expected_ext != ext:
            raise HTTPException(
                status_code=400,
                detail={
                    "error": "INVALID_FILE_TYPE",
                    "message": f"File '{filename}': MIME type '{content_type}' does not match extension '{ext}'.",
                },
            )

        max_bytes = settings.MAX_FILE_SIZE_MB * 1024 * 1024
        if size_bytes > max_bytes:
            raise HTTPException(
                status_code=400,
                detail={
                    "error": "FILE_TOO_LARGE",
                    "message": f"File '{filename}' is {size_bytes / (1024 * 1024):.1f} MB, "
                    f"which exceeds the {settings.MAX_FILE_SIZE_MB} MB limit.",
                },
            )

    def parse(self, filename: str, content: bytes, content_type: str) -> ParsedDocument:
        ext = os.path.splitext(filename)[1].lower()

        dispatch = {
            ".pdf": self._parse_pdf,
            ".docx": self._parse_docx,
            ".xlsx": self._parse_xlsx,
            ".pptx": self._parse_pptx,
            ".txt": self._parse_text,
            ".md": self._parse_text,
        }
        parser = dispatch.get(ext)
        if parser is None:
            raise HTTPException(
                status_code=400,
                detail={
                    "error": "INVALID_FILE_TYPE",
                    "message": f"Unsupported extension '{ext}'.",
                },
            )

        raw_text = parser(content)

        if len(raw_text.strip()) < 50:
            raise HTTPException(
                status_code=422,
                detail={
                    "error": "NO_TEXT_EXTRACTED",
                    "message": f"File '{filename}' yielded fewer than 50 characters of text. "
                    "Ensure the document contains readable content.",
                },
            )

        file_type = ext.lstrip(".")
        return ParsedDocument(
            filename=filename,
            file_type=file_type,
            raw_text=raw_text,
            char_count=len(raw_text),
        )

    def _parse_pdf(self, content: bytes) -> str:
        pages: list[str] = []
        with fitz.open(stream=content, filetype="pdf") as doc:
            for page in doc:
                pages.append(page.get_text())
        return "\n".join(pages)

    def _parse_docx(self, content: bytes) -> str:
        doc = DocxDocument(io.BytesIO(content))
        return "\n".join(p.text for p in doc.paragraphs if p.text)

    def _parse_xlsx(self, content: bytes) -> str:
        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        lines: list[str] = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                cells = [str(cell) for cell in row if cell is not None and str(cell).strip()]
                if cells:
                    lines.append("\t".join(cells))
        wb.close()
        return "\n".join(lines)

    def _parse_pptx(self, content: bytes) -> str:
        prs = Presentation(io.BytesIO(content))
        lines: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)
        return "\n".join(lines)

    def _parse_text(self, content: bytes) -> str:
        try:
            return content.decode("utf-8")
        except UnicodeDecodeError:
            return content.decode("latin-1")


def validate_total_size(files: list[UploadFile], settings_obj=None) -> None:
    if settings_obj is None:
        settings_obj = settings
    total = sum(f.size for f in files if f.size is not None)
    max_bytes = settings_obj.MAX_TOTAL_SIZE_MB * 1024 * 1024
    if total > max_bytes:
        raise HTTPException(
            status_code=400,
            detail={
                "error": "TOTAL_SIZE_EXCEEDED",
                "message": f"Total upload size {total / (1024 * 1024):.1f} MB exceeds "
                f"the {settings_obj.MAX_TOTAL_SIZE_MB} MB limit.",
            },
        )
