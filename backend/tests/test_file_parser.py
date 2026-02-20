import io
from pathlib import Path
from unittest.mock import MagicMock

import pytest
from docx import Document as DocxDocument
from fastapi import HTTPException, UploadFile
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from app.services.file_parser import FileParser, ParsedDocument, validate_total_size

FIXTURES = Path(__file__).parent / "fixtures"

parser = FileParser()


# ---------------------------------------------------------------------------
# Parsing happy-path tests
# ---------------------------------------------------------------------------


def test_parse_txt() -> None:
    content = (FIXTURES / "sample.txt").read_bytes()
    result = parser.parse("sample.txt", content, "text/plain")
    assert isinstance(result, ParsedDocument)
    assert result.file_type == "txt"
    assert result.filename == "sample.txt"
    assert "parse" in result.raw_text.lower()
    assert result.char_count == len(result.raw_text)


def test_parse_md() -> None:
    content = (FIXTURES / "sample.md").read_bytes()
    result = parser.parse("sample.md", content, "text/markdown")
    assert result.file_type == "md"
    assert "markdown" in result.raw_text.lower()
    assert result.char_count >= 50


def test_parse_docx() -> None:
    buf = io.BytesIO()
    doc = DocxDocument()
    doc.add_paragraph("User Story: As a user I want to log in so that I can access my data.")
    doc.add_paragraph("NFR: The system shall respond within 200ms under normal load.")
    doc.save(buf)

    result = parser.parse("test.docx", buf.getvalue(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    assert result.file_type == "docx"
    assert "user story" in result.raw_text.lower()
    assert "nfr" in result.raw_text.lower()


def test_parse_xlsx() -> None:
    buf = io.BytesIO()
    wb = Workbook()
    ws = wb.active
    ws.title = "Requirements"
    ws.append(["ID", "Title", "Priority"])
    ws.append(["US-001", "Login feature", "High"])
    ws.append(["US-002", "Logout feature", "Medium"])
    wb.save(buf)

    result = parser.parse("test.xlsx", buf.getvalue(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
    assert result.file_type == "xlsx"
    assert "Login feature" in result.raw_text
    assert "Logout feature" in result.raw_text
    # Three rows should each appear as a tab-separated line
    lines = [line for line in result.raw_text.splitlines() if line.strip()]
    assert len(lines) == 3


def test_parse_pptx() -> None:
    buf = io.BytesIO()
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Project Requirements Overview"
    slide.placeholders[1].text = "The system must support PDF, DOCX, XLSX, and PPTX files."
    prs.save(buf)

    result = parser.parse("test.pptx", buf.getvalue(), "application/vnd.openxmlformats-officedocument.presentationml.presentation")
    assert result.file_type == "pptx"
    assert "Requirements" in result.raw_text


# ---------------------------------------------------------------------------
# Validation error tests
# ---------------------------------------------------------------------------


def test_validate_unsupported_extension() -> None:
    with pytest.raises(HTTPException) as exc_info:
        parser.validate_file("malware.exe", "application/octet-stream", 1024)
    assert exc_info.value.status_code == 400
    assert exc_info.value.detail["error"] == "INVALID_FILE_TYPE"


def test_validate_file_too_large() -> None:
    size_31mb = 31 * 1024 * 1024
    with pytest.raises(HTTPException) as exc_info:
        parser.validate_file("big.pdf", "application/pdf", size_31mb)
    assert exc_info.value.status_code == 400
    assert exc_info.value.detail["error"] == "FILE_TOO_LARGE"
    assert "big.pdf" in exc_info.value.detail["message"]


def test_validate_total_size_exceeded() -> None:
    def make_upload(size: int) -> UploadFile:
        f = MagicMock(spec=UploadFile)
        f.size = size
        return f

    size_26mb = 26 * 1024 * 1024
    files = [make_upload(size_26mb), make_upload(size_26mb)]  # 52 MB total > 50 MB limit

    with pytest.raises(HTTPException) as exc_info:
        validate_total_size(files)
    assert exc_info.value.status_code == 400
    assert exc_info.value.detail["error"] == "TOTAL_SIZE_EXCEEDED"


def test_empty_document_raises() -> None:
    # A .txt file with only whitespace should raise NO_TEXT_EXTRACTED
    whitespace_content = b"   \n\t\n   "
    with pytest.raises(HTTPException) as exc_info:
        parser.parse("empty.txt", whitespace_content, "text/plain")
    assert exc_info.value.status_code == 422
    assert exc_info.value.detail["error"] == "NO_TEXT_EXTRACTED"
